# app.py - MultiRAG backend (LangChain v0.2+ compatible, Groq LLM, OCR support)
import os
import io
import time
import base64
import shutil
import logging
import threading
import uuid
import re
from typing import List, Dict, Any

from flask import Flask, request, jsonify, render_template, send_from_directory
from flask_cors import CORS
from werkzeug.utils import secure_filename

# LangChain/compat imports
from langchain_chroma import Chroma
try:
    from langchain_core.schema import Document
    from langchain_core.schema.messages import HumanMessage, SystemMessage
except Exception:
    from langchain.schema import Document
    from langchain.schema.messages import HumanMessage, SystemMessage

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_groq import ChatGroq
from dotenv import load_dotenv

from PIL import Image
import fitz  # PyMuPDF
import pdfplumber
import docx
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
import cv2

# --------------------------------------------------------------------------
# Setup
# --------------------------------------------------------------------------
load_dotenv()
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__, template_folder="templates")
CORS(app)

# Config
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # 200MB
app.config['UPLOAD_FOLDER'] = os.getenv('UPLOAD_FOLDER', 'uploads')
app.config['PROCESSED_FOLDER'] = os.getenv('PROCESSED_FOLDER', 'processed')
app.config['CHROMA_DIR'] = os.getenv('CHROMA_DIR', os.path.join(app.config['PROCESSED_FOLDER'], "chroma_db"))
app.config['ALLOWED_EXTENSIONS'] = set(os.getenv(
    'ALLOWED_EXTENSIONS',
    'pdf,doc,docx,xlsx,xls,ppt,pptx,txt,md,html,rtf,odt,ods,odp,csv,png,jpg,jpeg'
).split(','))
app.config['CHUNK_SIZE'] = int(os.getenv('CHUNK_SIZE', 1000))
app.config['CHUNK_OVERLAP'] = int(os.getenv('CHUNK_OVERLAP', 200))
app.config['EMBEDDING_MODEL'] = os.getenv('EMBEDDING_MODEL', 'sentence-transformers/all-MiniLM-L6-v2')
app.config['EMBEDDING_DEVICE'] = os.getenv('EMBEDDING_DEVICE', 'cpu')
app.config['SIMILARITY_K'] = int(os.getenv('SIMILARITY_K', 10))
app.config['SIMILARITY_THRESHOLD'] = float(os.getenv('SIMILARITY_THRESHOLD', 0.5))

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)
os.makedirs(app.config['CHROMA_DIR'], exist_ok=True)

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=app.config['CHUNK_SIZE'],
    chunk_overlap=app.config['CHUNK_OVERLAP'],
    separators=["\n\n", "\n", " ", ""]
)
embeddings = HuggingFaceEmbeddings(
    model_name=app.config['EMBEDDING_MODEL'],
    model_kwargs={'device': app.config['EMBEDDING_DEVICE']},
    encode_kwargs={'normalize_embeddings': True}
)
groq_client = ChatGroq(
    groq_api_key=os.getenv("GROQ_API_KEY"),
    model_name=os.getenv("GROQ_MODEL", "llama-3.1-8b-instant"),
    temperature=float(os.getenv("GROQ_TEMPERATURE", "0.1"))
)

vector_store: Chroma = None
vector_lock = threading.Lock()
all_docs: List[Document] = []
image_data_store: Dict[str, str] = {}
current_file_info: Dict[str, Any] = {}
upload_status: Dict[str, Dict[str, Any]] = {}
upload_status_lock = threading.Lock()
current_document_source: str = None

# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def allowed_file(filename: str) -> bool:
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def save_file_streamed(file_storage, dest_path: str, chunk_size: int = 16 * 1024):
    with open(dest_path, "wb") as w:
        shutil.copyfileobj(file_storage.stream, w, length=chunk_size)


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s.,!?;:-]', ' ', text)
    return text.strip()


def create_enhanced_context_message(query: str, context_docs: List[Document]) -> List:
    system_prompt = (
        "Answer the question using only the following context. "
        "If you don't know the answer from the context, reply: 'I don't know.'\nCONTEXT:\n"
    )
    context_block = "\n".join([clean_text(d.page_content) for d in context_docs]) if context_docs else "No relevant context found."
    return [SystemMessage(content=system_prompt + context_block), HumanMessage(content=query)]


def fast_clear_directory(dirpath: str):
    if os.path.exists(dirpath):
        shutil.rmtree(dirpath, ignore_errors=True)
    os.makedirs(dirpath, exist_ok=True)


def extract_text_from_image(path: str) -> str:
    try:
        img = cv2.imread(path)
        if img is None:
            return ""
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.adaptiveThreshold(gray, 255,
                                     cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                     cv2.THRESH_BINARY, 11, 2)
        text = pytesseract.image_to_string(gray, lang='eng')
        return text
    except Exception:
        return ""


# --------------------------------------------------------------------------
# Document Processor
# --------------------------------------------------------------------------
class AdvancedDocumentProcessor:
    """Extract text & images from multiple file formats."""

    def process_document(self, path: str, filename: str):
        ext = filename.lower().rsplit(".", 1)[-1]
        docs, images = [], {}

        if ext in ["pdf"]:
            docs, images = self._process_pdf(path, filename)
        elif ext in ["docx", "doc", "odt", "rtf"]:
            docs = self._process_docx(path, filename)
        elif ext in ["xlsx", "xls", "ods"]:
            docs = self._process_excel(path, filename)
        elif ext in ["pptx", "ppt", "odp"]:
            docs = self._process_ppt(path, filename)
        elif ext in ["txt", "md", "csv", "html"]:
            docs = self._process_text(path, filename)
        elif ext in ["png", "jpg", "jpeg"]:
            docs, images = self._process_image(path, filename)
        else:
            docs = [Document(page_content="Unsupported file format", metadata={"source": filename})]

        return docs, images

    # PDF
    def _process_pdf(self, path: str, filename: str):
        docs, images = [], {}
        try:
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        chunks = text_splitter.split_text(clean_text(text))
                        docs.extend([Document(page_content=c, metadata={"source": filename, "page": i + 1}) for c in chunks])
        except Exception:
            pass

        try:
            pdf_doc = fitz.open(path)
            for page_num in range(len(pdf_doc)):
                for img_index, img in enumerate(pdf_doc.get_page_images(page_num)):
                    xref = img[0]
                    base = pdf_doc.extract_image(xref)
                    image_bytes = base.get("image") if base else None
                    if image_bytes:
                        tmp_path = f"tmp_{uuid.uuid4().hex}.png"
                        with open(tmp_path, "wb") as f:
                            f.write(image_bytes)
                        text = extract_text_from_image(tmp_path)
                        os.remove(tmp_path)
                        if text.strip():
                            chunks = text_splitter.split_text(clean_text(text))
                            docs.extend([Document(page_content=c, metadata={"source": filename, "page": page_num + 1}) for c in chunks])
                        img_id = f"{uuid.uuid4().hex}"
                        images[img_id] = base64.b64encode(image_bytes).decode("utf-8")
        except Exception:
            pass

        return docs, images

    # DOCX/ODT
    def _process_docx(self, path: str, filename: str):
        docs = []
        try:
            doc = docx.Document(path)
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if full_text.strip():
                chunks = text_splitter.split_text(clean_text(full_text))
                docs.extend([Document(page_content=c, metadata={"source": filename}) for c in chunks])
        except Exception:
            pass
        return docs

    # Excel
    def _process_excel(self, path: str, filename: str):
        docs = []
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            for sheet in wb.sheetnames:
                sh = wb[sheet]
                values = []
                for row in sh.iter_rows(values_only=True):
                    vals = [str(cell) for cell in (row or []) if cell is not None and str(cell).strip()]
                    if vals:
                        values.append(" ".join(vals))
                text = "\n".join(values)
                if text.strip():
                    chunks = text_splitter.split_text(clean_text(text))
                    docs.extend([Document(page_content=c, metadata={"source": filename, "sheet": sheet}) for c in chunks])
        except Exception:
            pass
        return docs

    # PowerPoint
    def _process_ppt(self, path: str, filename: str):
        docs = []
        try:
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                text = "\n".join(texts)
                if text.strip():
                    chunks = text_splitter.split_text(clean_text(text))
                    docs.extend([Document(page_content=c, metadata={"source": filename, "slide": i + 1}) for c in chunks])
        except Exception:
            pass
        return docs

    # Text
    def _process_text(self, path: str, filename: str):
        docs = []
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            if text.strip():
                chunks = text_splitter.split_text(clean_text(text))
                docs.extend([Document(page_content=c, metadata={"source": filename}) for c in chunks])
        except Exception:
            pass
        return docs

    # Images
    def _process_image(self, path: str, filename: str):
        docs, images = [], {}
        try:
            text = extract_text_from_image(path)
            if text.strip():
                chunks = text_splitter.split_text(clean_text(text))
                docs.extend([Document(page_content=c, metadata={"source": filename}) for c in chunks])
            with open(path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")
            img_id = f"{uuid.uuid4().hex}"
            images[img_id] = b64
        except Exception:
            pass
        return docs, images


doc_processor = AdvancedDocumentProcessor()

# --------------------------------------------------------------------------
# Vector Store Handling
# --------------------------------------------------------------------------
def get_or_create_vector_store(docs: List[Document] = None) -> Chroma:
    global vector_store, current_document_source
    with vector_lock:
        if vector_store is None:
            try:
                vector_store = Chroma(persist_directory=app.config['CHROMA_DIR'], embedding_function=embeddings)
            except Exception:
                vector_store = None

        if docs and current_document_source:
            try:
                if vector_store:
                    all_ids = vector_store.get().get('ids', [])
                    if all_ids:
                        vector_store.delete(ids=all_ids)
            except Exception:
                vector_store = None

        if vector_store is None and docs:
            try:
                vector_store = Chroma.from_documents(documents=docs, embedding_function=embeddings, persist_directory=app.config['CHROMA_DIR'])
            except Exception as e:
                logger.exception("Failed to create Chroma: %s", e)
                vector_store = None
        elif vector_store and docs:
            try:
                vector_store.add_documents(docs)
                vector_store.persist()
            except Exception:
                pass
    return vector_store

# --------------------------------------------------------------------------
# Background Upload Processing
# --------------------------------------------------------------------------
def process_upload_background(upload_id: str, path: str, filename: str):
    global current_document_source
    with upload_status_lock:
        upload_status[upload_id]['status'] = 'processing'
        upload_status[upload_id]['step'] = 'extracting'
    try:
        docs, images = doc_processor.process_document(path, filename)
        with upload_status_lock:
            upload_status[upload_id]['chunks'] = len(docs)
            upload_status[upload_id]['images'] = len(images)
            upload_status[upload_id]['step'] = 'saving_images'
        for img_id, b64 in images.items():
            image_data_store[img_id] = b64
            try:
                outpath = os.path.join(app.config['PROCESSED_FOLDER'], f"{img_id}.png")
                with open(outpath, 'wb') as f:
                    f.write(base64.b64decode(b64))
            except Exception:
                pass
        current_document_source = filename
        with upload_status_lock:
            upload_status[upload_id]['step'] = 'indexing'
        vs = get_or_create_vector_store(docs)
        with upload_status_lock:
            upload_status[upload_id]['status'] = 'done' if vs else 'failed'
            upload_status[upload_id]['step'] = 'completed'
            upload_status[upload_id]['processed_at'] = time.time()
    except Exception as e:
        logger.exception("Processing upload failed: %s", e)
        with upload_status_lock:
            upload_status[upload_id]['status'] = 'failed'
            upload_status[upload_id]['error'] = str(e)

# --------------------------------------------------------------------------
# Routes
# --------------------------------------------------------------------------
@app.route('/')
def index():
    return render_template('index.html') if os.path.exists('templates/index.html') else jsonify({'status': 'ok', 'message': 'MultiRAG API'})


@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'status': 'error', 'message': 'No file provided'}), 400
    file = request.files['file']
    if not file or file.filename == '':
        return jsonify({'status': 'error', 'message': 'No file selected'}), 400
    if not allowed_file(file.filename):
        return jsonify({'status': 'error', 'message': f'File type not allowed: {file.filename}'}), 400

    filename = secure_filename(file.filename)
    upload_id = uuid.uuid4().hex
    dest_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{upload_id}_{filename}")
    try:
        save_file_streamed(file, dest_path)
    except Exception as e:
        logger.exception("Failed saving upload: %s", e)
        return jsonify({'status': 'error', 'message': f'Failed to save file: {e}'}), 500

    with upload_status_lock:
        upload_status[upload_id] = {
            'status': 'queued',
            'filename': filename,
            'chunks': 0,
            'images': 0,
            'step': 'queued',
            'error': None,
            'created_at': time.time()
        }

    threading.Thread(target=process_upload_background, args=(upload_id, dest_path, filename), daemon=True).start()
    return jsonify({'status': 'success', 'upload_id': upload_id, 'message': f'{filename} uploaded and processing started'}), 202


@app.route('/upload_status/<upload_id>', methods=['GET'])
def upload_status_route(upload_id: str):
    with upload_status_lock:
        st = upload_status.get(upload_id)
    if not st:
        return jsonify({'status': 'error', 'message': 'upload_id not found'}), 404
    resp = dict(st)
    if resp.get('processed_at'):
        resp['processed_at_human'] = time.ctime(resp['processed_at'])
    return jsonify(resp)


@app.route('/ask', methods=['POST'])
def ask():
    payload = request.get_json(force=True)
    query = payload.get('query', '').strip()
    if not query:
        return jsonify({'status': 'error', 'message': 'Query required'}), 400
    vs = get_or_create_vector_store()
    if vs is None:
        return jsonify({'status': 'error', 'message': 'No vector DB. Upload documents first.'}), 400
    filter_dict = {"source": current_document_source} if current_document_source else None
    results = []
    try:
        results = vs.similarity_search_with_score(query, k=app.config['SIMILARITY_K'] * 2, filter=filter_dict)
    except Exception:
        try:
            results = [(d, 1.0) for d in vs.similarity_search(query, k=app.config['SIMILARITY_K'] * 2)]
        except Exception as e:
            logger.exception("Similarity search failed: %s", e)
            return jsonify({'status': 'error', 'message': 'Similarity search failed.'}), 500

    filtered_results = [doc for doc, score in results if score is not None and score >= app.config['SIMILARITY_THRESHOLD']]
    top_docs = filtered_results[:app.config['SIMILARITY_K']]

    messages = create_enhanced_context_message(query, top_docs) if top_docs else [
        SystemMessage(content="Answer from the uploaded file context. If not found, say: 'I don't know.'"),
        HumanMessage(content=query)
    ]

    try:
        response = groq_client.invoke(messages)
        llm_answer = getattr(response, 'content', str(response))
    except Exception as e:
        logger.exception("LLM invocation error: %s", e)
        return jsonify({'status': 'error', 'message': f'LLM error: {e}'}), 500

    return jsonify({
        'status': 'success',
        'answer': llm_answer,
        'retrieved_documents': [{'source': d.metadata.get('source'), 'preview': d.page_content[:200]} for d in top_docs]
    })


@app.route('/clear', methods=['POST'])
@app.route('/clear_data', methods=['POST'])
def clear_data():
    global vector_store, all_docs, image_data_store, current_file_info, upload_status, current_document_source
    with vector_lock:
        vector_store = None
    all_docs = []
    image_data_store = {}
    current_file_info = {}
    current_document_source = None
    with upload_status_lock:
        upload_status = {}
    fast_clear_directory(app.config['PROCESSED_FOLDER'])
    fast_clear_directory(app.config['CHROMA_DIR'])
    return jsonify({'status': 'success', 'message': 'Cleared vector DB and files.'})


@app.route('/get_image/<image_id>')
def get_image(image_id: str):
    filename = f"{image_id}.png"
    path = os.path.join(app.config['PROCESSED_FOLDER'], filename)
    if os.path.exists(path):
        return send_from_directory(app.config['PROCESSED_FOLDER'], filename)
    b64 = image_data_store.get(image_id)
    if b64:
        return (base64.b64decode(b64), 200, {'Content-Type': 'image/png'})
    return jsonify({'status': 'error', 'message': 'Image not found'}), 404


@app.route('/list_images')
def list_images():
    try:
        imgs = [fn.replace('.png', '') for fn in os.listdir(app.config['PROCESSED_FOLDER']) if fn.lower().endswith('.png')]
    except Exception:
        imgs = list(image_data_store.keys())
    return jsonify({'status': 'success', 'images': imgs})


@app.route('/health')
def health():
    return jsonify({
        'status': 'ok',
        'has_vector_store': bool(vector_store),
        'loaded_docs': len(all_docs),
        'current_document': current_document_source
    })

# --------------------------------------------------------------------------
# Main
# --------------------------------------------------------------------------
if __name__ == '__main__':
    logger.info("Starting MultiRAG server. Chroma DB present: %s", os.path.exists(app.config['CHROMA_DIR']))
    app.run(debug=True, host='0.0.0.0', port=int(os.getenv('PORT', 5000)))
